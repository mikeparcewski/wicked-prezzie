#!/usr/bin/env python3
"""
PPTX Builder — Converts extracted layout data into editable PowerPoint slides.

Takes the JSON output from chrome_extract and creates native PPTX shapes,
text boxes with inline formatting, and embedded images at correct positions.
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Color utilities (bundled in this skill)
from color_utils import parse_css_color, is_bold, decode_entities

FONT = "Calibri"


def pp_align(css_align):
    """Convert CSS text-align to python-pptx alignment."""
    if css_align == 'center': return PP_ALIGN.CENTER
    if css_align == 'right': return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


class SlideBuilder:
    """
    Builds a PPTX slide from extracted layout data.

    Coordinate system: layout data uses pixel coordinates in the HTML viewport
    space (default 1280x720). These are converted to PPTX EMU units.
    """

    def __init__(self, slide_w_inches=13.333, slide_h_inches=7.5,
                 source_w=1280, source_h=720, font=FONT):
        self.sw = slide_w_inches
        self.sh = slide_h_inches
        self.src_w = source_w
        self.src_h = source_h
        self.font = font

    def px2emu_x(self, px):
        return int(px / self.src_w * self.sw * 914400)

    def px2emu_y(self, px):
        return int(px / self.src_h * self.sh * 914400)

    def build_slide(self, prs, layout_data, screenshot_fn=None, notes_text=None):
        """
        Build a single PPTX slide from layout data.

        Args:
            prs: python-pptx Presentation object
            layout_data: dict from chrome_extract.extract_layout()
            screenshot_fn: callable(svg_rect) -> png_path for SVG rendering
            notes_text: optional speaker notes text

        Returns:
            The created slide object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Background
        bg_color = parse_css_color(layout_data.get('slideBg'))
        slide_cls = layout_data.get('slideClasses', [])
        if 'section-divider' in slide_cls:
            if not bg_color or bg_color == RGBColor(0, 0, 0):
                bg_color = RGBColor(0x14, 0x00, 0x24)
        elif 'title-slide' in slide_cls:
            if not bg_color:
                bg_color = RGBColor(0, 0, 0)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color or RGBColor(0x0A, 0x0A, 0x0F)

        slide_bg_rgb = (10, 10, 15)
        if bg_color:
            slide_bg_rgb = (bg_color[0], bg_color[1], bg_color[2])

        elements = layout_data.get('elements', [])
        svg_elements = layout_data.get('svgElements', [])

        shapes_data = sorted([e for e in elements if e['type'] == 'shape'],
                            key=lambda e: e['depth'])

        # Filter shapes
        filtered_shapes = []
        for s in shapes_data:
            r = s['rect']
            if r['w'] > 1250 and r['h'] > 680: continue
            if r['w'] < 4 or r['h'] < 4: continue
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            if r['x'] + r['w'] < 0 or r['y'] + r['h'] < 0: continue
            filtered_shapes.append(s)

        # Helper: find parent card shape for text clamping
        def find_parent_card(text_rect):
            best = None
            best_area = float('inf')
            for s in filtered_shapes:
                sr = s['rect']
                if (text_rect['x'] >= sr['x'] - 5 and text_rect['y'] >= sr['y'] - 5 and
                    text_rect['x'] + text_rect['w'] <= sr['x'] + sr['w'] + 30 and
                    text_rect['y'] + text_rect['h'] <= sr['y'] + sr['h'] + 10):
                    area = sr['w'] * sr['h']
                    if area < best_area and sr['w'] > 20 and sr['h'] > 20:
                        best = sr
                        best_area = area
            return best

        # Render shapes
        for s in filtered_shapes:
            r = s['rect']
            st = s['styles']
            x = self.px2emu_x(max(0, r['x']))
            y = self.px2emu_y(max(0, r['y']))
            w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
            h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))
            if w < 10000 or h < 8000: continue

            br = st.get('borderRadius', 0)
            if br > 5:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
                try:
                    shape.adjustments[0] = min(0.12, br / min(r['w'], r['h']))
                except: pass
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

            fill_c = parse_css_color(st.get('backgroundColor'), slide_bg_rgb)
            if fill_c:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_c
            else:
                shape.fill.background()

            border_c = parse_css_color(st.get('borderColor'), slide_bg_rgb)
            bw = st.get('borderWidth', 0)
            blw = st.get('borderLeftWidth', 0)
            blc = parse_css_color(st.get('borderLeftColor'), slide_bg_rgb)
            if blc and blw > 2:
                shape.line.color.rgb = blc
                shape.line.width = Pt(max(1, min(blw, 4)))
            elif border_c and bw > 0.3:
                shape.line.color.rgb = border_c
                shape.line.width = Pt(max(0.5, min(bw, 3)))
            else:
                shape.line.fill.background()
            shape.text_frame.clear()

        # Render SVGs as images
        if screenshot_fn:
            for svg in svg_elements:
                r = svg['rect']
                if r['w'] < 60 or r['h'] < 60: continue
                if svg.get('lines', 0) < 10: continue
                png_path = screenshot_fn(r)
                if png_path and os.path.exists(png_path):
                    x = self.px2emu_x(max(0, r['x']))
                    y = self.px2emu_y(max(0, r['y']))
                    w = self.px2emu_x(r['w'])
                    h = self.px2emu_y(r['h'])
                    try:
                        slide.shapes.add_picture(png_path, x, y, w, h)
                    except: pass

        # Separate richtext and simple text
        richtext_data = [e for e in elements if e['type'] == 'richtext']
        simple_text_data = sorted([e for e in elements if e['type'] == 'text'],
                                  key=lambda e: e['depth'])

        richtext_rects = [rt['rect'] for rt in richtext_data]

        def is_covered_by_richtext(r):
            for rr in richtext_rects:
                if (r['x'] >= rr['x'] - 5 and r['y'] >= rr['y'] - 5 and
                    r['x'] + r['w'] <= rr['x'] + rr['w'] + 10 and
                    r['y'] + r['h'] <= rr['y'] + rr['h'] + 10):
                    return True
            return False

        filtered_texts = []
        for t in simple_text_data:
            r = t['rect']
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            text = t.get('text', '').strip()
            if not text: continue
            if is_covered_by_richtext(r): continue
            dup = False
            for ex in filtered_texts:
                if (ex['text'].strip() == text and
                    abs(ex['rect']['x'] - r['x']) < 15 and
                    abs(ex['rect']['y'] - r['y']) < 15):
                    if t['depth'] > ex['depth']:
                        filtered_texts.remove(ex)
                    else:
                        dup = True
                    break
            if not dup:
                filtered_texts.append(t)

        # Render richtext elements
        for rt in richtext_data:
            r = rt['rect']
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            runs = rt.get('runs', [])
            if not runs: continue
            full_text = ''.join(run.get('text', '') for run in runs).strip()
            if not full_text: continue

            align = rt.get('styles', {}).get('textAlign', 'left')
            tag = rt.get('tag', '')
            parent = find_parent_card(r)

            if parent and tag not in ('h1',):
                px = parent['x'] + 4
                pw = parent['w'] - 8
                x = self.px2emu_x(max(0, px))
                y = self.px2emu_y(max(0, r['y']))
                w = self.px2emu_x(max(20, pw))
            elif align == 'center' and tag in ('h1', 'h2', 'h3', 'p'):
                x = self.px2emu_x(48)
                y = self.px2emu_y(max(0, r['y']))
                w = self.px2emu_x(self.src_w - 96)
            else:
                x = self.px2emu_x(max(0, r['x']))
                y = self.px2emu_y(max(0, r['y']))
                extra = 1.15 if tag in ('h1', 'h2', 'h3') else 1.08
                w = self.px2emu_x(min(r['w'] * extra, self.src_w - max(0, r['x'])))

            h = self.px2emu_y(min(max(r['h'], 14), self.src_h - max(0, r['y'])))
            if w < 5000 or h < 5000: continue

            tb = slide.shapes.add_textbox(x, y, w, h)
            tb.fill.background()
            tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

            p = tf.paragraphs[0]
            p.alignment = pp_align(align)
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            for run_data in runs:
                text = run_data.get('text', '')
                if not text: continue
                if run_data.get('br'):
                    p = tf.add_paragraph()
                    p.alignment = pp_align(align)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    continue
                if '\n' in text:
                    parts = text.split('\n')
                    for pi, part in enumerate(parts):
                        if pi > 0:
                            p = tf.add_paragraph()
                            p.alignment = pp_align(align)
                            p.space_before = Pt(0)
                            p.space_after = Pt(0)
                        if part.strip():
                            self._add_run(p, part, run_data)
                    continue
                self._add_run(p, text, run_data)

        # Render simple text
        for t in filtered_texts:
            r = t['rect']
            st = t['styles']
            text = t.get('text', '').strip()
            if not text: continue

            tt = st.get('textTransform', '')
            if tt == 'uppercase':
                text = text.upper()

            parent = find_parent_card(r)
            if parent:
                px = parent['x'] + 4
                pw = parent['w'] - 8
                x = self.px2emu_x(max(0, px))
                w = self.px2emu_x(max(20, pw))
            else:
                x = self.px2emu_x(max(0, r['x']))
                w = self.px2emu_x(min(r['w'] * 1.02, self.src_w - max(0, r['x'])))
            y = self.px2emu_y(max(0, r['y']))
            h = self.px2emu_y(min(max(r['h'], 14), self.src_h - max(0, r['y'])))
            if w < 5000 or h < 5000: continue

            tb = slide.shapes.add_textbox(x, y, w, h)
            tb.fill.background()
            tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

            p = tf.paragraphs[0]
            p.alignment = pp_align(st.get('textAlign', 'left'))
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            run = p.add_run()
            run.text = decode_entities(text)
            run.font.name = self.font
            fs = st.get('fontSize', 14)
            run.font.size = Pt(round(max(6, min(52, fs * 0.75)), 1))
            run.font.bold = is_bold(st.get('fontWeight'))
            run.font.italic = st.get('fontStyle') == 'italic'
            color = parse_css_color(st.get('color'), (255, 255, 255))
            if color:
                run.font.color.rgb = color

        # Speaker notes
        if notes_text:
            try:
                slide.notes_slide.notes_text_frame.text = notes_text
            except: pass

        return slide

    def _add_run(self, paragraph, text, run_data):
        """Add a formatted run to a paragraph."""
        tt = run_data.get('textTransform', '')
        if tt == 'uppercase':
            text = text.upper()
        run = paragraph.add_run()
        run.text = decode_entities(text)
        run.font.name = self.font
        fs = run_data.get('fontSize', 14)
        run.font.size = Pt(round(max(6, min(52, fs * 0.75)), 1))
        run.font.bold = is_bold(run_data.get('fontWeight'))
        run.font.italic = run_data.get('fontStyle') == 'italic'
        color = parse_css_color(run_data.get('color'), (255, 255, 255))
        if color:
            run.font.color.rgb = color
        return run
