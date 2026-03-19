#!/usr/bin/env python3
"""
PPTX Builder — Converts extracted layout data into editable PowerPoint slides.

Takes the JSON output from chrome_extract and creates native PPTX shapes,
text boxes with inline formatting, and embedded images at correct positions.
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Color utilities (bundled in this skill)
from color_utils import parse_css_color, is_bold, decode_entities

FONT = "Calibri"
EMOJI_FONT = "Segoe UI Emoji"  # Windows/cross-platform; macOS falls back to Apple Color Emoji

import platform as _platform
if _platform.system() == 'Darwin':
    EMOJI_FONT = "Apple Color Emoji"


def _has_emoji(text):
    """Check if text contains emoji characters that Calibri can't render."""
    for ch in text:
        cp = ord(ch)
        if cp >= 0x1F300:  # Emoji range (miscellaneous symbols and pictographs+)
            return True
        if 0x2600 <= cp <= 0x27BF:  # Misc symbols, dingbats
            return True
        if 0xFE00 <= cp <= 0xFE0F:  # Variation selectors
            return True
        if 0x200D == cp:  # Zero-width joiner (emoji sequences)
            return True
    return False


def pp_align(css_align):
    """Convert CSS text-align to python-pptx alignment."""
    if css_align == 'center': return PP_ALIGN.CENTER
    if css_align == 'right': return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


import re as _re
import math as _math

def _parse_css_gradient(grad_str):
    """Parse a CSS linear-gradient string into angle and color stops."""
    if not grad_str or 'linear-gradient' not in str(grad_str):
        return None, []
    angle_match = _re.search(r'(\d+)deg', grad_str)
    angle = float(angle_match.group(1)) if angle_match else 180.0
    stops = _re.findall(r'rgba?\(([^)]+)\)', grad_str)
    parsed = []
    for s in stops:
        parts = [x.strip() for x in s.split(',')]
        r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
        a = float(parts[3]) if len(parts) == 4 else 1.0
        parsed.append((r, g, b, a))
    return angle, parsed


def _apply_gradient_fill(shape, angle_css, stops, bg_rgb):
    """Apply OOXML gradient fill if stops differ enough."""
    blended = []
    for r, g, b, a in stops:
        br, bg_g, bb = bg_rgb
        blended.append((
            int(r * a + br * (1 - a)),
            int(g * a + bg_g * (1 - a)),
            int(b * a + bb * (1 - a))
        ))
    if len(blended) < 2:
        return False
    dist = _math.sqrt(sum((a - b) ** 2 for a, b in zip(blended[0], blended[-1])))
    if dist < 20:
        return False  # imperceptible gradient, use solid
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill') or child.tag.endswith('}noFill'):
            spPr.remove(child)
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    for i, (cr, cg, cb) in enumerate(blended):
        pos = int(i / max(1, len(blended) - 1) * 100000)
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgbClr = etree.SubElement(gs, qn('a:srgbClr'))
        srgbClr.set('val', f'{cr:02X}{cg:02X}{cb:02X}')
    lin = etree.SubElement(gradFill, qn('a:lin'))
    ooxml_deg = (90 - angle_css) % 360
    lin.set('ang', str(int(ooxml_deg * 60000)))
    lin.set('scaled', '0')
    return True


class SlideBuilder:
    """
    Builds a PPTX slide from extracted layout data.

    Coordinate system: layout data uses pixel coordinates in the HTML viewport
    space (default 1280x720). These are converted to PPTX EMU units.
    """

    def __init__(self, slide_w_inches=13.333, slide_h_inches=7.5,
                 source_w=1280, source_h=720, font=FONT):
        self.sw = slide_w_inches
        self.sh = slide_h_inches
        self.src_w = source_w
        self.src_h = source_h
        self.font = font

    def px2emu_x(self, px):
        return int(px / self.src_w * self.sw * 914400)

    def px2emu_y(self, px):
        return int(px / self.src_h * self.sh * 914400)

    def _apply_margins(self, tf, styles=None):
        """Apply CSS padding as text frame margins, or zero margins if no padding."""
        PX_TO_EMU = 9525  # 1px at 96 DPI
        if styles and any(styles.get(k, 0) for k in ('paddingTop', 'paddingRight', 'paddingBottom', 'paddingLeft')):
            tf.margin_top = Emu(int(styles.get('paddingTop', 0) * PX_TO_EMU))
            tf.margin_bottom = Emu(int(styles.get('paddingBottom', 0) * PX_TO_EMU))
            tf.margin_left = Emu(int(styles.get('paddingLeft', 0) * PX_TO_EMU))
            tf.margin_right = Emu(int(styles.get('paddingRight', 0) * PX_TO_EMU))
        else:
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    def build_slide(self, prs, layout_data, screenshot_fn=None, notes_text=None):
        """
        Build a single PPTX slide from layout data.

        Args:
            prs: python-pptx Presentation object
            layout_data: dict from chrome_extract.extract_layout()
            screenshot_fn: callable(svg_rect) -> png_path for SVG rendering
            notes_text: optional speaker notes text

        Returns:
            The created slide object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Background
        bg_color = parse_css_color(layout_data.get('slideBg'))
        slide_cls = layout_data.get('slideClasses', [])
        if 'section-divider' in slide_cls:
            if not bg_color or bg_color == RGBColor(0, 0, 0):
                bg_color = RGBColor(0x14, 0x00, 0x24)
        elif 'title-slide' in slide_cls:
            if not bg_color:
                bg_color = RGBColor(0, 0, 0)
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color or RGBColor(0x0A, 0x0A, 0x0F)

        slide_bg_rgb = (10, 10, 15)
        if bg_color:
            slide_bg_rgb = (bg_color[0], bg_color[1], bg_color[2])

        elements = layout_data.get('elements', [])
        svg_elements = layout_data.get('svgElements', [])

        shapes_data = sorted([e for e in elements if e['type'] == 'shape'],
                            key=lambda e: e['depth'])

        # Filter shapes
        filtered_shapes = []
        for s in shapes_data:
            r = s['rect']
            if r['w'] > 1250 and r['h'] > 680: continue
            if r['w'] * r['h'] < 16: continue  # truly tiny
            if r['w'] < 2 and r['h'] < 2: continue
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            if r['x'] + r['w'] < 0 or r['y'] + r['h'] < 0: continue
            filtered_shapes.append(s)

        # Helper: find parent card shape for text clamping
        def find_parent_card(text_rect):
            best = None
            best_area = float('inf')
            for s in filtered_shapes:
                sr = s['rect']
                if (text_rect['x'] >= sr['x'] - 5 and text_rect['y'] >= sr['y'] - 5 and
                    text_rect['x'] + text_rect['w'] <= sr['x'] + sr['w'] + 30 and
                    text_rect['y'] + text_rect['h'] <= sr['y'] + sr['h'] + 10):
                    area = sr['w'] * sr['h']
                    if area < best_area and sr['w'] > 20 and sr['h'] > 20:
                        best = sr
                        best_area = area
            return best

        def find_containment_rect(text_node):
            """Tightest ancestor rect: explicit flex/grid slot, or card ancestor."""
            if text_node.get('parentSlotRect'):
                return text_node['parentSlotRect']
            return find_parent_card(text_node['rect'])

        # Render shapes
        for s in filtered_shapes:
            r = s['rect']
            st = s['styles']
            x = self.px2emu_x(max(0, r['x']))
            y = self.px2emu_y(max(0, r['y']))
            w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
            h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))
            if w < 10000 and h < 8000: continue  # allow thin decorative lines

            br = st.get('borderRadius', 0)
            if br > 5:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
                try:
                    shape.adjustments[0] = min(0.12, br / min(r['w'], r['h']))
                except: pass
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

            fill_c = parse_css_color(st.get('backgroundColor'), slide_bg_rgb)
            if fill_c:
                shape.fill.solid()
                shape.fill.fore_color.rgb = fill_c
            else:
                shape.fill.background()

            # Try gradient fill from raw CSS background
            bg_raw = st.get('background', '')
            if bg_raw and 'linear-gradient' in str(bg_raw):
                angle, stops = _parse_css_gradient(bg_raw)
                if stops:
                    _apply_gradient_fill(shape, angle, stops, slide_bg_rgb)

            border_c = parse_css_color(st.get('borderColor'), slide_bg_rgb)
            bw = st.get('borderWidth', 0)
            blw = st.get('borderLeftWidth', 0)
            blc = parse_css_color(st.get('borderLeftColor'), slide_bg_rgb)
            if blc and blw > 2:
                # Create separate accent bar instead of full outline (#25)
                accent_w = self.px2emu_x(max(3, min(blw, 6)))
                accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, accent_w, h)
                accent.fill.solid()
                accent.fill.fore_color.rgb = blc
                accent.line.fill.background()
                accent.text_frame.clear()
                shape.line.fill.background()
            elif border_c and bw > 0.3:
                shape.line.color.rgb = border_c
                shape.line.width = Pt(max(0.25, bw * 0.75))  # CSS px → PPTX pt at 96 DPI
            else:
                shape.line.fill.background()
            shape.text_frame.clear()

        # Render SVGs as images
        if screenshot_fn:
            # Collect non-SVG element top edges for overlap detection (#19)
            non_svg_tops = []
            for e in elements:
                if e['type'] in ('richtext', 'text', 'shape', 'badge'):
                    non_svg_tops.append(e['rect']['y'])

            for svg in svg_elements:
                r = svg['rect']
                if r['w'] < 30 or r['h'] < 30: continue
                if svg.get('lines', 0) < 3: continue
                # Clamp SVG bottom to avoid bleeding into content below (#19)
                # Elements within 30px below SVG can bleed through transparent areas
                svg_bottom = r['y'] + r['h']
                crop_r = dict(r)
                nearest_below = None
                for ey in sorted(non_svg_tops):
                    if ey > r['y'] + r['h'] * 0.5 and ey <= svg_bottom + 30:
                        nearest_below = ey
                        break
                if nearest_below is not None:
                    new_h = nearest_below - r['y'] - 8
                    if new_h > r['h'] * 0.5 and new_h < r['h']:
                        crop_r['h'] = new_h
                png_path = screenshot_fn(crop_r)
                if png_path and os.path.exists(png_path):
                    x = self.px2emu_x(max(0, crop_r['x']))
                    y = self.px2emu_y(max(0, crop_r['y']))
                    w = self.px2emu_x(crop_r['w'])
                    h = self.px2emu_y(crop_r['h'])
                    try:
                        slide.shapes.add_picture(png_path, x, y, w, h)
                    except: pass

        # Render native tables
        table_data = [e for e in elements if e['type'] == 'table'
                      and e['rect']['x'] < self.src_w and e['rect']['y'] < self.src_h
                      and e['rect']['x'] + e['rect']['w'] > 0 and e['rect']['y'] + e['rect']['h'] > 0]
        for td in table_data:
            try:
                self.build_table(slide, td, slide_bg_rgb)
            except Exception:
                pass

        # Render badges and circles
        badge_data = [e for e in elements if e['type'] in ('badge', 'circle')]
        for bd in badge_data:
            try:
                self.build_badge(slide, bd, slide_bg_rgb)
            except Exception:
                pass

        # Separate richtext and simple text
        richtext_data = [e for e in elements if e['type'] == 'richtext']
        simple_text_data = sorted([e for e in elements if e['type'] == 'text'],
                                  key=lambda e: e['depth'])

        richtext_rects = [rt['rect'] for rt in richtext_data]

        def is_covered_by_richtext(r):
            for rr in richtext_rects:
                if (r['x'] >= rr['x'] - 5 and r['y'] >= rr['y'] - 5 and
                    r['x'] + r['w'] <= rr['x'] + rr['w'] + 10 and
                    r['y'] + r['h'] <= rr['y'] + rr['h'] + 10):
                    return True
            return False

        filtered_texts = []
        for t in simple_text_data:
            r = t['rect']
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            text = t.get('text', '').strip()
            if not text: continue
            if is_covered_by_richtext(r): continue
            dup = False
            for ex in filtered_texts:
                if (ex['text'].strip() == text and
                    abs(ex['rect']['x'] - r['x']) < 15 and
                    abs(ex['rect']['y'] - r['y']) < 15):
                    if t['depth'] > ex['depth']:
                        filtered_texts.remove(ex)
                    else:
                        dup = True
                    break
            if not dup:
                filtered_texts.append(t)

        # Render richtext elements
        for rt in richtext_data:
            r = rt['rect']
            if r['x'] > self.src_w or r['y'] > self.src_h: continue
            runs = rt.get('runs', [])
            if not runs: continue
            full_text = ''.join(run.get('text', '') for run in runs).strip()
            if not full_text: continue

            rotation = rt.get('rotation', 0)
            align = rt.get('styles', {}).get('textAlign', 'left')
            tag = rt.get('tag', '')
            parent = find_containment_rect(rt)

            if parent and tag not in ('h1',):
                px = parent['x'] + 4
                pw = parent['w'] - 8
                x = self.px2emu_x(max(0, px))
                y = self.px2emu_y(max(0, r['y']))
                # Card width is a hard ceiling (#29 bug 7)
                w = self.px2emu_x(max(20, min(pw, parent['w'])))
            elif align == 'center' and tag in ('h1', 'h2', 'h3', 'p'):
                # Center within actual width, not forced to full-slide (#29 bug 7)
                actual_w = min(r['w'] * 1.05, self.src_w - 96)
                cx = (self.src_w - actual_w) / 2
                x = self.px2emu_x(max(0, cx))
                y = self.px2emu_y(max(0, r['y']))
                w = self.px2emu_x(actual_w)
            else:
                x = self.px2emu_x(max(0, r['x']))
                y = self.px2emu_y(max(0, r['y']))
                extra = 1.05 if tag in ('h1', 'h2', 'h3') else 1.03
                w = self.px2emu_x(min(r['w'] * extra, self.src_w - max(0, r['x'])))

            h = self.px2emu_y(min(max(r['h'], 14), self.src_h - max(0, r['y'])))
            if w < 5000 or h < 5000: continue

            # For rotated text, swap dimensions and apply PPTX rotation (#27)
            if rotation and abs(rotation) >= 5:
                w, h = h, w
            tb = slide.shapes.add_textbox(x, y, w, h)
            if rotation and abs(rotation) >= 5:
                tb.rotation = rotation
            tb.fill.background()
            tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = True
            ws = rt.get('styles', {}).get('whiteSpace', '')
            if ws == 'nowrap':
                tf.word_wrap = False
            self._apply_margins(tf, rt.get('styles', {}))

            p = tf.paragraphs[0]
            p.alignment = pp_align(align)
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            for run_data in runs:
                text = run_data.get('text', '')
                if not text: continue
                if run_data.get('br'):
                    p = tf.add_paragraph()
                    p.alignment = pp_align(align)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    continue
                if '\n' in text:
                    parts = text.split('\n')
                    for pi, part in enumerate(parts):
                        if pi > 0:
                            p = tf.add_paragraph()
                            p.alignment = pp_align(align)
                            p.space_before = Pt(0)
                            p.space_after = Pt(0)
                        if part.strip():
                            self._add_run(p, part, run_data)
                    continue
                self._add_run(p, text, run_data)

        # Render simple text
        for t in filtered_texts:
            r = t['rect']
            st = t['styles']
            text = t.get('text', '').strip()
            if not text: continue

            rotation = t.get('rotation', 0)
            tt = st.get('textTransform', '')
            if tt == 'uppercase':
                text = text.upper()

            parent = find_containment_rect(t)
            if parent:
                px = parent['x'] + 4
                pw = parent['w'] - 8
                x = self.px2emu_x(max(0, px))
                w = self.px2emu_x(max(20, min(pw, parent['w'])))
            else:
                x = self.px2emu_x(max(0, r['x']))
                w = self.px2emu_x(min(r['w'] * 1.02, self.src_w - max(0, r['x'])))
            y = self.px2emu_y(max(0, r['y']))
            h = self.px2emu_y(min(max(r['h'], 14), self.src_h - max(0, r['y'])))
            if w < 5000 or h < 5000: continue

            # For rotated text, swap dimensions (#27)
            if rotation and abs(rotation) >= 5:
                w, h = h, w
            tb = slide.shapes.add_textbox(x, y, w, h)
            if rotation and abs(rotation) >= 5:
                tb.rotation = rotation
            tb.fill.background()
            tb.line.fill.background()
            tf = tb.text_frame
            tf.word_wrap = True
            self._apply_margins(tf, st)

            p = tf.paragraphs[0]
            p.alignment = pp_align(st.get('textAlign', 'left'))
            p.space_before = Pt(0)
            p.space_after = Pt(0)

            run = p.add_run()
            decoded_text = decode_entities(text)
            run.text = decoded_text
            run.font.name = EMOJI_FONT if _has_emoji(decoded_text) else self.font
            fs = st.get('fontSize', 14)
            run.font.size = Pt(round(max(6, min(52, fs * 0.75)), 1))
            run.font.bold = is_bold(st.get('fontWeight'))
            run.font.italic = st.get('fontStyle') == 'italic'
            color = parse_css_color(st.get('color'), (255, 255, 255))
            if color:
                run.font.color.rgb = color

        # Speaker notes
        if notes_text:
            try:
                slide.notes_slide.notes_text_frame.text = notes_text
            except: pass

        return slide

    def _add_run(self, paragraph, text, run_data):
        """Add a formatted run to a paragraph."""
        tt = run_data.get('textTransform', '')
        if tt == 'uppercase':
            text = text.upper()
        run = paragraph.add_run()
        run.text = decode_entities(text)
        run.font.name = EMOJI_FONT if _has_emoji(text) else self.font
        fs = run_data.get('fontSize', 14)
        run.font.size = Pt(round(max(6, min(52, fs * 0.75)), 1))
        run.font.bold = is_bold(run_data.get('fontWeight'))
        run.font.italic = run_data.get('fontStyle') == 'italic'
        ls = run_data.get('letterSpacing')
        if ls and ls != 'normal':
            try:
                px = float(ls) if isinstance(ls, (int, float)) else float(str(ls).replace('px', ''))
                if abs(px) >= 0.1:
                    rPr = run._r.get_or_add_rPr()
                    rPr.set('spc', str(int(px * 75)))  # px → 100ths of pt
            except (ValueError, TypeError):
                pass
        color = parse_css_color(run_data.get('color'), (255, 255, 255))
        if color:
            run.font.color.rgb = color
        return run

    def build_table(self, slide, table_node, slide_bg_rgb):
        """Build a native PPTX table from extracted table data."""
        rows = table_node.get('rows', [])
        if not rows:
            return None
        num_rows = len(rows)
        num_cols = max(len(row) for row in rows) if rows else 0
        if num_cols == 0:
            return None
        r = table_node['rect']
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))

        table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, w, h)
        table = table_shape.table

        # Disable default table styling that causes white banded rows (#16)
        tbl = table._tbl
        tblPr = tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            tblPr.attrib.clear()
            tblPr.set('bandRow', '0')
            tblPr.set('bandCol', '0')
            tblPr.set('firstRow', '0')
            tblPr.set('lastRow', '0')
            tblPr.set('firstCol', '0')
            tblPr.set('lastCol', '0')
            for child in list(tblPr):
                if child.tag.endswith('}tblStyle'):
                    tblPr.remove(child)

        # Set column widths from first row cell widths
        if rows[0]:
            for ci, cell_data in enumerate(rows[0]):
                if ci < num_cols:
                    col_w = self.px2emu_x(cell_data['rect']['w'])
                    tbl.tblGrid[ci].set('w', str(col_w))

        fallback_bg = RGBColor(slide_bg_rgb[0], slide_bg_rgb[1], slide_bg_rgb[2])

        for ri, row in enumerate(rows):
            for ci, cell_data in enumerate(row):
                if ci >= num_cols or ri >= num_rows:
                    continue
                cell = table.cell(ri, ci)
                bg = parse_css_color(cell_data['styles'].get('backgroundColor'), slide_bg_rgb)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg or fallback_bg
                tf = cell.text_frame
                tf.word_wrap = True
                runs = cell_data.get('runs', [])
                first_run = True
                for run_data in runs:
                    if run_data.get('br'):
                        tf.add_paragraph()
                        continue
                    text = run_data.get('text', '').strip()
                    if not text:
                        continue
                    if first_run:
                        p = tf.paragraphs[0]
                        first_run = False
                    else:
                        p = tf.paragraphs[-1]
                    run = p.add_run()
                    run.text = text
                    run.font.name = self.font
                    run.font.size = Pt(round(max(6, min(52, run_data.get('fontSize', 12) * 0.75)), 1))
                    color = parse_css_color(run_data.get('color'), (255, 255, 255))
                    if color:
                        run.font.color.rgb = color
                    if is_bold(run_data.get('fontWeight')):
                        run.font.bold = True
        return table_shape

    def build_badge(self, slide, node, slide_bg_rgb):
        """Build a pill/circle badge shape."""
        r = node['rect']
        # Bounds check: skip badges outside the slide viewport (#17)
        if r['x'] > self.src_w or r['y'] > self.src_h:
            return None
        if r['x'] + r['w'] < 0 or r['y'] + r['h'] < 0:
            return None
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))
        st = node.get('styles', {})
        br = st.get('borderRadius', 0)
        if br >= min(r['w'], r['h']) / 2:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            try:
                shape.adjustments[0] = 0.5  # max rounding = pill
            except Exception:
                pass
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            try:
                shape.adjustments[0] = min(0.15, br / max(1, min(r['w'], r['h'])))
            except Exception:
                pass
        bg = parse_css_color(st.get('backgroundColor'), slide_bg_rgb)
        if bg:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
        else:
            shape.fill.background()
        shape.line.fill.background()
        text = node.get('text', '').strip()
        if text:
            tf = shape.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = text
            run.font.name = self.font
            run.font.size = Pt(round(max(6, st.get('fontSize', 12) * 0.75), 1))
            color = parse_css_color(st.get('color'), (255, 255, 255))
            if color:
                run.font.color.rgb = color
