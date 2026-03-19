#!/usr/bin/env python3
"""EDL (Edit Description Language) applicator for PPTX slides.

Reads a JSON edit spec and applies non-destructive edits to an existing PPTX file.
This is the safe execution surface for model-driven fixes — the model writes the
JSON spec, this script applies it. No python-pptx API hallucination possible.

Usage:
    python edl_apply.py deck.pptx edits.json -o deck-fixed.pptx

EDL Schema:
    {
      "edits": [
        {
          "slide": 3,                    # 1-based slide index
          "target": "shape:5",           # shape index, or "text:substring", or "image:0"
          "ops": [
            {"action": "move", "x": 100, "y": 200},
            {"action": "resize", "width": 400, "height": 300},
            {"action": "crop_bottom", "pixels": 30},
            {"action": "set_text", "text": "New text"},
            {"action": "set_fill", "color": "#A100FF"},
            {"action": "set_font_size", "size": 14},
            {"action": "delete"},
            {"action": "add_shape", "shape": "rectangle",
             "x": 10, "y": 100, "width": 4, "height": 80,
             "fill": "#22c55e"}
          ]
        }
      ]
    }
"""

import argparse
import json
import os
import sys
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Import sibling color utils
_root = Path(__file__).parent
sys.path.insert(0, str(_root))
from color_utils import parse_css_color


# Coordinate helpers (match pptx_builder.py defaults)
DEFAULT_SW = 13.333
DEFAULT_SH = 7.5
DEFAULT_SRC_W = 1280
DEFAULT_SRC_H = 720


def px2emu_x(px, sw=DEFAULT_SW, src_w=DEFAULT_SRC_W):
    return int(px / src_w * sw * 914400)


def px2emu_y(px, sh=DEFAULT_SH, src_h=DEFAULT_SRC_H):
    return int(px / src_h * sh * 914400)


def hex_to_rgb(hex_str):
    """Convert #RRGGBB to RGBColor."""
    h = hex_str.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def find_target(slide, target_spec):
    """Resolve a target spec to a shape object.

    Formats:
        "shape:N"       — shape by index (0-based)
        "text:substring" — first shape containing substring
        "image:N"       — Nth image shape (0-based)
    """
    if target_spec.startswith('shape:'):
        idx = int(target_spec.split(':')[1])
        shapes = list(slide.shapes)
        if 0 <= idx < len(shapes):
            return shapes[idx]
    elif target_spec.startswith('text:'):
        needle = target_spec[5:]
        for shape in slide.shapes:
            if shape.has_text_frame and needle in shape.text_frame.text:
                return shape
    elif target_spec.startswith('image:'):
        idx = int(target_spec.split(':')[1])
        images = [s for s in slide.shapes if s.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
        if 0 <= idx < len(images):
            return images[idx]
    return None


def apply_op(slide, shape, op):
    """Apply a single operation to a shape."""
    action = op['action']

    if action == 'move':
        if 'x' in op:
            shape.left = px2emu_x(op['x'])
        if 'y' in op:
            shape.top = px2emu_y(op['y'])

    elif action == 'resize':
        if 'width' in op:
            shape.width = px2emu_x(op['width'])
        if 'height' in op:
            shape.height = px2emu_y(op['height'])

    elif action == 'crop_bottom':
        px = op.get('pixels', 0)
        shape.height -= px2emu_y(px)

    elif action == 'crop_top':
        px = op.get('pixels', 0)
        shape.top += px2emu_y(px)
        shape.height -= px2emu_y(px)

    elif action == 'set_text' and shape.has_text_frame:
        shape.text_frame.paragraphs[0].runs[0].text = op['text']

    elif action == 'set_fill':
        color = hex_to_rgb(op['color'])
        shape.fill.solid()
        shape.fill.fore_color.rgb = color

    elif action == 'set_font_size' and shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.size = Pt(op['size'])

    elif action == 'set_rotation':
        shape.rotation = op.get('degrees', 0)

    elif action == 'delete':
        sp = shape._element
        sp.getparent().remove(sp)

    elif action == 'add_shape':
        # This creates a new shape on the slide, ignoring `shape` param
        x = px2emu_x(op.get('x', 0))
        y = px2emu_y(op.get('y', 0))
        w = px2emu_x(op.get('width', 100))
        h = px2emu_y(op.get('height', 50))
        shape_type = {
            'rectangle': MSO_SHAPE.RECTANGLE,
            'rounded_rectangle': MSO_SHAPE.ROUNDED_RECTANGLE,
            'oval': MSO_SHAPE.OVAL,
        }.get(op.get('shape', 'rectangle'), MSO_SHAPE.RECTANGLE)
        new_shape = slide.shapes.add_shape(shape_type, x, y, w, h)
        if 'fill' in op:
            new_shape.fill.solid()
            new_shape.fill.fore_color.rgb = hex_to_rgb(op['fill'])
        else:
            new_shape.fill.background()
        new_shape.line.fill.background()

    elif action == 'add_textbox':
        x = px2emu_x(op.get('x', 0))
        y = px2emu_y(op.get('y', 0))
        w = px2emu_x(op.get('width', 200))
        h = px2emu_y(op.get('height', 30))
        tb = slide.shapes.add_textbox(x, y, w, h)
        tb.fill.background()
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        run = tf.paragraphs[0].add_run()
        run.text = op.get('text', '')
        run.font.name = op.get('font', 'Calibri')
        run.font.size = Pt(op.get('font_size', 14))
        if 'color' in op:
            run.font.color.rgb = hex_to_rgb(op['color'])

    return True


def apply_edl(pptx_path, edl, output_path=None):
    """Apply an EDL edit spec to a PPTX file.

    Args:
        pptx_path: Path to input PPTX
        edl: dict with 'edits' key (parsed JSON)
        output_path: Output path (defaults to overwriting input)

    Returns:
        dict with results per edit
    """
    prs = Presentation(pptx_path)
    results = []

    for edit in edl.get('edits', []):
        slide_idx = edit.get('slide', 1) - 1  # 1-based to 0-based
        if slide_idx < 0 or slide_idx >= len(prs.slides):
            results.append({'slide': edit.get('slide'), 'status': 'error', 'reason': 'slide out of range'})
            continue

        slide = prs.slides[slide_idx]
        target_spec = edit.get('target', '')
        ops = edit.get('ops', [])

        for op in ops:
            if op['action'] in ('add_shape', 'add_textbox'):
                # These don't need a target
                apply_op(slide, None, op)
                results.append({'slide': edit.get('slide'), 'action': op['action'], 'status': 'ok'})
            else:
                shape = find_target(slide, target_spec)
                if shape is None:
                    results.append({'slide': edit.get('slide'), 'target': target_spec, 'status': 'error', 'reason': 'target not found'})
                    continue
                apply_op(slide, shape, op)
                results.append({'slide': edit.get('slide'), 'target': target_spec, 'action': op['action'], 'status': 'ok'})

    out = output_path or pptx_path
    prs.save(out)
    return {'output': str(out), 'results': results}


def main():
    parser = argparse.ArgumentParser(description='Apply EDL edits to a PPTX file.')
    parser.add_argument('pptx', help='Input PPTX file')
    parser.add_argument('edl', help='EDL JSON file with edit specs')
    parser.add_argument('-o', '--output', default=None, help='Output PPTX path (default: overwrite input)')
    args = parser.parse_args()

    with open(args.edl) as f:
        edl = json.load(f)

    result = apply_edl(args.pptx, edl, args.output)
    print(json.dumps(result, indent=2))


if __name__ == '__main__':
    main()
