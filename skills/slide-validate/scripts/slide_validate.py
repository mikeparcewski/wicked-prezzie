#!/usr/bin/env python3
"""Slide Validate — Structural checks for PowerPoint files.

Fast pre-check that catches shape overflow, negative coordinates, empty slides,
and text overflow heuristics. Visual fidelity is judged by Claude comparing
HTML screenshots vs PPTX renders — not by this script.

Usage from other skills:
    from slide_validate import validate_pptx
    report = validate_pptx("deck.pptx")
"""

import os, json
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu


def check_bounds(slide, slide_w, slide_h):
    """Check shapes for overflow beyond slide boundaries."""
    issues = []
    tolerance = Emu(50000)
    for shape in slide.shapes:
        right = shape.left + shape.width
        bottom = shape.top + shape.height
        if right > slide_w + tolerance:
            overflow = right - slide_w
            issues.append({
                "type": "bounds", "severity": "error",
                "description": f"Shape '{shape.name}' extends {Emu(overflow).inches:.2f}in past right edge",
            })
        if bottom > slide_h + tolerance:
            overflow = bottom - slide_h
            issues.append({
                "type": "bounds", "severity": "error",
                "description": f"Shape '{shape.name}' extends {Emu(overflow).inches:.2f}in past bottom edge",
            })
    return issues


def check_negative_coords(slide):
    """Detect shapes with negative positions (partially off-slide)."""
    issues = []
    tolerance = Emu(-50000)
    for shape in slide.shapes:
        if shape.left < tolerance:
            issues.append({
                "type": "negative_coord", "severity": "warning",
                "description": f"Shape '{shape.name}' has negative left position",
            })
        if shape.top < tolerance:
            issues.append({
                "type": "negative_coord", "severity": "warning",
                "description": f"Shape '{shape.name}' has negative top position",
            })
    return issues


def check_empty_slide(slide):
    """Detect slides with no content."""
    if len(slide.shapes) == 0:
        return [{"type": "empty", "severity": "warning",
                 "description": "Slide has no shapes"}]
    has_text = any(
        shape.has_text_frame and shape.text_frame.text.strip()
        for shape in slide.shapes if hasattr(shape, "has_text_frame")
    )
    has_image = any(shape.shape_type == 13 for shape in slide.shapes)
    if not has_text and not has_image:
        return [{"type": "empty", "severity": "warning",
                 "description": "Slide has no text or images"}]
    return []


def check_text_overflow(slide):
    """Heuristic check for text that likely overflows its container."""
    issues = []
    for shape in slide.shapes:
        if not hasattr(shape, "has_text_frame") or not shape.has_text_frame:
            continue
        text = shape.text_frame.text
        if not text.strip():
            continue
        width_inches = shape.width / 914400
        estimated_chars_per_line = width_inches * 12
        estimated_lines = shape.height / 914400 / 0.25
        max_chars = estimated_chars_per_line * max(1, estimated_lines)
        if len(text) > max_chars * 1.5:
            issues.append({
                "type": "text_overflow", "severity": "warning",
                "description": (
                    f"Text in '{shape.name}' may overflow "
                    f"({len(text)} chars in {width_inches:.1f}in wide box)"
                ),
            })
    return issues


def validate_pptx(pptx_path):
    """Run structural validation on a PPTX file.

    Returns:
        List of per-slide dicts with index and issues list.
    """
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    results = []

    for i, slide in enumerate(prs.slides):
        issues = []
        issues.extend(check_bounds(slide, slide_w, slide_h))
        issues.extend(check_negative_coords(slide))
        issues.extend(check_empty_slide(slide))
        issues.extend(check_text_overflow(slide))
        results.append({"index": i + 1, "issues": issues})

    return results
