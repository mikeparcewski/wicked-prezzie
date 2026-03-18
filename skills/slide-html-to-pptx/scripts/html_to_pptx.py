#!/usr/bin/env python3
"""
slide-html-to-pptx skill — Convert HTML slide decks to editable PowerPoint.

Parallel extraction: each slide is standardized, extracted, and screenshotted
concurrently via ProcessPoolExecutor. Results are assembled into one PPTX
sequentially (python-pptx is not thread-safe).

Usage:
    python skills/html_to_pptx.py --input-dir ./slides --output deck.pptx
    python skills/html_to_pptx.py --slides slide-01.html,slide-02.html --output deck.pptx
    python skills/html_to_pptx.py --manifest slides.json --output deck.pptx
"""

import argparse, os, sys, json, tempfile, glob, time
from pathlib import Path
from concurrent.futures import ProcessPoolExecutor, as_completed

# Import from sibling skills
_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(_root / "shared"))
from paths import output_path
sys.path.insert(0, str(_root / "chrome-extract" / "scripts"))
sys.path.insert(0, str(_root / "slide-pptx-builder" / "scripts"))
sys.path.insert(0, str(_root / "slide-html-standardize" / "scripts"))

from chrome_extract import extract_layout, screenshot_html, crop_region
from pptx_builder import SlideBuilder
from html_standardize import standardize_html
from pptx import Presentation
from pptx.util import Inches
from bs4 import BeautifulSoup


def extract_notes(html_path):
    """Extract text from HTML for speaker notes."""
    try:
        with open(html_path, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f.read(), 'lxml')
        for tag in soup.find_all(['script', 'style', 'nav']):
            tag.decompose()
        sl = soup.find(class_='slide') or soup.body
        return sl.get_text(separator='\n', strip=True)[:3000] if sl else ''
    except:
        return ''


# ---------------------------------------------------------------------------
# Phase 1: Per-slide extraction (runs in worker processes)
# ---------------------------------------------------------------------------

def _extract_single_slide(args):
    """
    Worker function: standardize + extract layout + take screenshot for one slide.

    Runs in a separate process. Returns a dict with all data needed to build
    the PPTX slide in the assembly phase. The screenshot is taken once here
    and cached on disk — no duplicate Chrome launches during assembly.
    """
    (index, slide_info, input_dir, cache_dir,
     viewport_w, viewport_h, hide_selectors, do_standardize) = args

    filename = slide_info['file']
    html_path = os.path.join(input_dir, filename)

    result = {
        'index': index,
        'slide_info': slide_info,
        'filename': filename,
        'layout': None,
        'screenshot_path': None,
        'notes': '',
        'error': None,
    }

    if not os.path.exists(html_path):
        result['error'] = 'NOT FOUND'
        return result

    try:
        # Standardize if requested
        if do_standardize:
            standardize_html(html_path, viewport_w=viewport_w, viewport_h=viewport_h)

        # Per-slide temp directory for Chrome artifacts
        slide_tmpdir = os.path.join(cache_dir, f"slide-{index:03d}")
        os.makedirs(slide_tmpdir, exist_ok=True)

        # Extract layout via Chrome headless
        layout = extract_layout(html_path, slide_tmpdir, viewport_w, viewport_h,
                                hide_selectors)
        result['layout'] = layout

        # Take the full-page screenshot once — used for SVG cropping and fallback
        ss_path = os.path.join(cache_dir, f"screenshot-{index:03d}.png")
        screenshot_html(html_path, ss_path, slide_tmpdir,
                        viewport_w, viewport_h, hide_selectors=hide_selectors)
        if os.path.exists(ss_path):
            result['screenshot_path'] = ss_path

        # Extract speaker notes
        result['notes'] = extract_notes(html_path)

    except Exception as e:
        result['error'] = str(e)

    return result


def _parallel_extract(work_items, max_workers, total):
    """Run parallel Chrome extraction and return results keyed by index."""
    extracted = {}
    with ProcessPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(_extract_single_slide, item): item[0]
                   for item in work_items}
        for future in as_completed(futures):
            result = future.result()
            idx = result['index']
            extracted[idx] = result
            filename = result['filename']
            if result['error'] == 'NOT FOUND':
                print(f"    [{idx+1}/{total}] {filename}: NOT FOUND")
            elif result['error']:
                print(f"    [{idx+1}/{total}] {filename}: error: {result['error']}")
            elif result['layout'] and result['layout'].get('elements'):
                ne = len(result['layout']['elements'])
                print(f"    [{idx+1}/{total}] {filename}: {ne} elements")
            else:
                print(f"    [{idx+1}/{total}] {filename}: fallback")
    return extracted


# ---------------------------------------------------------------------------
# Phase 2: PPTX assembly (runs sequentially in main process)
# ---------------------------------------------------------------------------

def build_deck(slides, input_dir, output_path, hide_selectors=None,
               viewport_w=1280, viewport_h=720, standardize=False,
               workers=None):
    """
    Build a PowerPoint deck from HTML slides using parallel extraction.

    Phase 1: Extract all slides concurrently (Chrome headless per slide).
    Phase 2: Assemble layouts into one PPTX sequentially.
    """
    sw, sh = 13.333, 7.5
    builder = SlideBuilder(sw, sh, viewport_w, viewport_h)
    prs = Presentation()
    prs.slide_width = Inches(sw)
    prs.slide_height = Inches(sh)

    total = len(slides)
    hide = hide_selectors or ['.slide-nav']
    # Cap default workers to avoid Chrome resource contention
    cpu = os.cpu_count() or 4
    max_workers = workers or min(total, cpu // 2 or 1, 6)

    with tempfile.TemporaryDirectory() as cache_dir:
        # --- Phase 1: Parallel extraction ---
        t0 = time.time()
        print(f"\n  Extracting {total} slides ({max_workers} workers)...")

        work_items = [
            (i, si, input_dir, cache_dir, viewport_w, viewport_h, hide, standardize)
            for i, si in enumerate(slides)
        ]

        extracted = _parallel_extract(work_items, max_workers, total)

        # --- Phase 1b: Retry failed extractions sequentially ---
        failed_indices = [
            idx for idx, r in extracted.items()
            if r['error'] and r['error'] != 'NOT FOUND'
        ]
        if failed_indices:
            print(f"\n  Retrying {len(failed_indices)} failed slide(s) sequentially...")
            for idx in failed_indices:
                item = work_items[idx]
                retry = _extract_single_slide(item)
                filename = retry['filename']
                if retry['error']:
                    print(f"    [{idx+1}/{total}] {filename}: retry failed: {retry['error']}")
                else:
                    ne = len(retry['layout'].get('elements', [])) if retry['layout'] else 0
                    print(f"    [{idx+1}/{total}] {filename}: retry OK ({ne} elements)")
                    extracted[idx] = retry

        extract_time = time.time() - t0
        print(f"  Extraction: {extract_time:.1f}s ({extract_time/total:.1f}s/slide effective)")

        # --- Phase 2: Sequential PPTX assembly ---
        t1 = time.time()
        print(f"\n  Assembling PPTX...")
        stats = {'layout': 0, 'fallback': 0, 'skipped': 0}

        for i in range(total):
            result = extracted.get(i)
            if not result or result['error'] == 'NOT FOUND':
                stats['skipped'] += 1
                continue

            si = result['slide_info']
            title = si.get('title', result['filename'])
            act = si.get('act', '')
            layout = result['layout']
            screenshot_path = result['screenshot_path']

            has_layout = (layout and 'error' not in layout
                          and len(layout.get('elements', [])) > 0)

            if has_layout:
                # SVG screenshot callback uses the cached full-page screenshot
                def screenshot_svg(svg_rect, _ss=screenshot_path, _idx=i,
                                   _cache=cache_dir, _vw=viewport_w, _vh=viewport_h):
                    if not _ss or not os.path.exists(_ss):
                        return None
                    out = os.path.join(_cache,
                        f"svg_{_idx}_{int(svg_rect['x'])}_{int(svg_rect['y'])}.png")
                    if crop_region(_ss, out, svg_rect, _vw, _vh):
                        return out
                    return None

                notes = f"[{act}] {title}\n\n{result['notes']}" if act else result['notes']
                try:
                    builder.build_slide(prs, layout, screenshot_svg, notes)
                    stats['layout'] += 1
                except Exception as e:
                    # Layout build failed — try fallback
                    _build_fallback_from_cache(prs, builder, screenshot_path,
                                               result['notes'])
                    stats['fallback'] += 1
            elif screenshot_path:
                # No usable layout — use cached screenshot as fallback
                _build_fallback_from_cache(prs, builder, screenshot_path,
                                           result['notes'])
                stats['fallback'] += 1
            elif not result['error']:
                # Extraction returned nothing and no screenshot
                stats['fallback'] += 1

        assemble_time = time.time() - t1

    prs.save(output_path)
    total_time = time.time() - t0
    print(f"\n  Saved: {output_path}")
    print(f"  Layout: {stats['layout']}, Fallback: {stats['fallback']}, "
          f"Skipped: {stats['skipped']}")
    print(f"  Timing: extract {extract_time:.1f}s + assemble {assemble_time:.1f}s "
          f"= {total_time:.1f}s total")
    return output_path


def _build_fallback_from_cache(prs, builder, screenshot_path, notes_text):
    """Build a fallback slide from a cached screenshot PNG."""
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0F)
    if screenshot_path and os.path.exists(screenshot_path):
        slide.shapes.add_picture(screenshot_path, Emu(0), Emu(0),
                                 Inches(builder.sw), Inches(builder.sh))
    if notes_text:
        slide.notes_slide.notes_text_frame.text = notes_text


# Keep the old fallback for backward compatibility with direct callers
def _build_fallback(prs, builder, html_path, si, tmpdir, vw, vh, hide_selectors):
    """Screenshot fallback when extraction fails (legacy sequential path)."""
    from pptx.dml.color import RGBColor
    from pptx.util import Emu
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0x0A, 0x0A, 0x0F)
    png = os.path.join(tmpdir, f"fb_{os.path.basename(html_path)}.png")
    if screenshot_html(html_path, png, tmpdir, vw, vh, hide_selectors=hide_selectors):
        slide.shapes.add_picture(png, Emu(0), Emu(0),
                                Inches(builder.sw), Inches(builder.sh))
    notes = extract_notes(html_path)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def main():
    parser = argparse.ArgumentParser(description='Convert HTML slides to editable PowerPoint')
    parser.add_argument('--input-dir', '-d', default='.', help='Directory containing HTML files')
    parser.add_argument('--output', '-o', default='deck.pptx', help='Output PPTX path')
    parser.add_argument('--slides', '-s', help='Comma-separated HTML filenames')
    parser.add_argument('--manifest', '-m', help='JSON manifest file with slide list')
    parser.add_argument('--hide', help='CSS selectors to hide (comma-separated)')
    parser.add_argument('--viewport', default='1280x720', help='Viewport WxH (default 1280x720)')
    parser.add_argument('--standardize', action='store_true',
                        help='Run slide-html-standardize on slides before conversion')
    parser.add_argument('--workers', '-w', type=int, default=None,
                        help='Max parallel workers (default: min(slide_count, cpu_count))')
    args = parser.parse_args()

    vw, vh = map(int, args.viewport.split('x'))
    hide = args.hide.split(',') if args.hide else ['.slide-nav']

    if args.manifest:
        with open(args.manifest) as f:
            slides = json.load(f)
    elif args.slides:
        slides = [{'file': f.strip()} for f in args.slides.split(',')]
    else:
        # Auto-discover HTML files
        html_files = sorted(glob.glob(os.path.join(args.input_dir, '*.html')))
        slides = [{'file': os.path.basename(f)} for f in html_files
                  if 'index' not in os.path.basename(f).lower()
                  and 'sorter' not in os.path.basename(f).lower()]

    if not slides:
        print("No slides found. Use --slides, --manifest, or place HTML files in --input-dir")
        sys.exit(1)

    print(f"Converting {len(slides)} slides from {args.input_dir}")
    build_deck(slides, args.input_dir, args.output, hide, vw, vh,
               standardize=args.standardize, workers=args.workers)


if __name__ == '__main__':
    main()
