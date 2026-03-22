#!/usr/bin/env python3
"""
PPTX Builder — Converts extracted layout data into editable PowerPoint slides.

Takes the JSON output from chrome_extract and creates native PPTX shapes,
text boxes with inline formatting, and embedded images at correct positions.
"""

import os, sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# Color utilities (bundled in this skill)
from color_utils import parse_css_color, is_bold, decode_entities

FONT = "Calibri"
EMOJI_FONT = "Segoe UI Emoji"  # Windows/cross-platform; macOS falls back to Apple Color Emoji

import platform as _platform
if _platform.system() == 'Darwin':
    EMOJI_FONT = "Apple Color Emoji"


def _has_emoji(text):
    """Check if text contains emoji characters that Calibri can't render."""
    for ch in text:
        cp = ord(ch)
        if cp >= 0x1F300:  # Emoji range (miscellaneous symbols and pictographs+)
            return True
        if 0x2600 <= cp <= 0x27BF:  # Misc symbols, dingbats
            return True
        if 0xFE00 <= cp <= 0xFE0F:  # Variation selectors
            return True
        if 0x200D == cp:  # Zero-width joiner (emoji sequences)
            return True
    return False


def _is_emoji_char(ch):
    """Check if a single character is a symbol/emoji that Calibri can't render well.

    Covers emoji, dingbats, geometric shapes, bullets, arrows, enclosed
    alphanumerics (①②③), checkmarks, and other decorative Unicode.
    """
    cp = ord(ch)
    if cp >= 0x1F300: return True               # Emoji (pictographs+)
    if 0x2600 <= cp <= 0x27BF: return True       # Misc symbols + dingbats
    if 0x2500 <= cp <= 0x25FF: return True       # Box drawing + geometric shapes (●▸◆)
    if 0x2B00 <= cp <= 0x2BFF: return True       # Misc symbols (⬢⬤)
    if 0xFE00 <= cp <= 0xFE0F: return True       # Variation selectors
    if cp == 0x200D: return True                 # Zero-width joiner
    if cp == 0x2022: return True                 # Bullet •
    if 0x2000 <= cp <= 0x206F: return True       # General punctuation (bullets •, dashes)
    if 0x2190 <= cp <= 0x21FF: return True       # Arrows (→←↑↓)
    if 0x2300 <= cp <= 0x23FF: return True       # Misc technical (⌘)
    if 0x2400 <= cp <= 0x24FF: return True       # Enclosed alphanumerics (①②③)
    if 0x2700 <= cp <= 0x27BF: return True       # Dingbats (✓✗✦)
    if 0x2010 <= cp <= 0x2027: return True       # Bullets and hyphens (•‣)
    if 0x2030 <= cp <= 0x205E: return True       # Per mille, primes, etc.
    if 0x2200 <= cp <= 0x22FF: return True       # Math operators (≈≠≤)
    if 0x2100 <= cp <= 0x214F: return True       # Letterlike symbols (™℃)
    return False


def _split_emoji(text):
    """Split text into (segment, is_emoji) tuples for font splitting.

    Adjacent emoji chars are grouped. Adjacent non-emoji chars are grouped.
    Returns list of (text, bool) tuples.
    """
    if not text:
        return []
    segments = []
    current = text[0]
    current_emoji = _is_emoji_char(text[0])
    for ch in text[1:]:
        ch_emoji = _is_emoji_char(ch)
        if ch_emoji == current_emoji:
            current += ch
        else:
            segments.append((current, current_emoji))
            current = ch
            current_emoji = ch_emoji
    segments.append((current, current_emoji))
    return segments


def pp_align(css_align):
    """Convert CSS text-align to python-pptx alignment."""
    if css_align == 'center': return PP_ALIGN.CENTER
    if css_align == 'right': return PP_ALIGN.RIGHT
    return PP_ALIGN.LEFT


import re as _re
import math as _math

def _parse_css_gradient(grad_str):
    """Parse a CSS linear-gradient string into angle and color stops."""
    if not grad_str or 'linear-gradient' not in str(grad_str):
        return None, []
    angle_match = _re.search(r'(\d+)deg', grad_str)
    angle = float(angle_match.group(1)) if angle_match else 180.0
    stops = _re.findall(r'rgba?\(([^)]+)\)', grad_str)
    parsed = []
    for s in stops:
        parts = [x.strip() for x in s.split(',')]
        r, g, b = int(parts[0]), int(parts[1]), int(parts[2])
        a = float(parts[3]) if len(parts) == 4 else 1.0
        parsed.append((r, g, b, a))
    return angle, parsed


def _apply_gradient_fill(shape, angle_css, stops, bg_rgb):
    """Apply OOXML gradient fill if stops differ enough."""
    blended = []
    for r, g, b, a in stops:
        br, bg_g, bb = bg_rgb
        blended.append((
            int(r * a + br * (1 - a)),
            int(g * a + bg_g * (1 - a)),
            int(b * a + bb * (1 - a))
        ))
    if len(blended) < 2:
        return False
    dist = _math.sqrt(sum((a - b) ** 2 for a, b in zip(blended[0], blended[-1])))
    if dist < 20:
        return False  # imperceptible gradient, use solid
    spPr = shape._element.spPr
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill') or child.tag.endswith('}noFill'):
            spPr.remove(child)
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    for i, (cr, cg, cb) in enumerate(blended):
        pos = int(i / max(1, len(blended) - 1) * 100000)
        gs = etree.SubElement(gsLst, qn('a:gs'))
        gs.set('pos', str(pos))
        srgbClr = etree.SubElement(gs, qn('a:srgbClr'))
        srgbClr.set('val', f'{cr:02X}{cg:02X}{cb:02X}')
    lin = etree.SubElement(gradFill, qn('a:lin'))
    ooxml_deg = (90 - angle_css) % 360
    lin.set('ang', str(int(ooxml_deg * 60000)))
    lin.set('scaled', '0')
    return True


class SlideBuilder:
    """
    Builds a PPTX slide from extracted layout data.

    Coordinate system: layout data uses pixel coordinates in the HTML viewport
    space (default 1280x720). These are converted to PPTX EMU units.
    """

    def __init__(self, slide_w_inches=13.333, slide_h_inches=7.5,
                 source_w=1280, source_h=720, font=FONT):
        self.sw = slide_w_inches
        self.sh = slide_h_inches
        self.src_w = source_w
        self.src_h = source_h
        self.font = font

    def px2emu_x(self, px):
        return int(px / self.src_w * self.sw * 914400)

    def px2emu_y(self, px):
        return int(px / self.src_h * self.sh * 914400)

    def _apply_margins(self, tf, styles=None):
        """Apply CSS padding as text frame margins, or zero margins if no padding."""
        PX_TO_EMU = 9525  # 1px at 96 DPI
        if styles and any(styles.get(k, 0) for k in ('paddingTop', 'paddingRight', 'paddingBottom', 'paddingLeft')):
            tf.margin_top = Emu(int(styles.get('paddingTop', 0) * PX_TO_EMU))
            tf.margin_bottom = Emu(int(styles.get('paddingBottom', 0) * PX_TO_EMU))
            tf.margin_left = Emu(int(styles.get('paddingLeft', 0) * PX_TO_EMU))
            tf.margin_right = Emu(int(styles.get('paddingRight', 0) * PX_TO_EMU))
        else:
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)

    def build_slide(self, prs, layout_data, screenshot_fn=None, notes_text=None):
        """
        Build a single PPTX slide from layout data (legacy compatibility wrapper).

        Converts layout_data to a build manifest via prep's auto_resolve(),
        then delegates to build_slide_from_manifest(). All classification logic
        and geometry transforms live in prep (known-patterns.md).

        Args:
            prs: python-pptx Presentation object
            layout_data: dict from chrome_extract.classify_elements()
            screenshot_fn: callable(svg_rect) -> png_path for SVG rendering
            notes_text: optional speaker notes text

        Returns:
            The created slide object
        """
        # Convert layout_data to a manifest via auto_resolve
        try:
            import sys as _sys
            _root = Path(__file__).parent.parent.parent
            if str(_root / "triage" / "scripts") not in _sys.path:
                _sys.path.insert(0, str(_root / "triage" / "scripts"))
            if str(_root / "prep" / "scripts") not in _sys.path:
                _sys.path.insert(0, str(_root / "prep" / "scripts"))
            from slide_triage import triage_slide
            from slide_prep import auto_resolve

            findings = triage_slide(raw_layout=layout_data,
                                    classified_layout=layout_data)
            manifest = auto_resolve(findings, layout_data)
        except Exception:
            # Fallback: build a minimal manifest from layout_data directly
            manifest = self._layout_to_manifest(layout_data)

        # Inject notes and slide class background overrides
        if notes_text:
            manifest['notes'] = notes_text
        slide_cls = layout_data.get('slideClasses', [])
        if not manifest.get('slideBg') or manifest.get('slideBg') == '#0A0A0F':
            if 'section-divider' in slide_cls:
                manifest['slideBg'] = '#140024'
            elif 'title-slide' in slide_cls:
                manifest['slideBg'] = '#000000'

        return self.build_slide_from_manifest(prs, manifest, screenshot_fn)

    def _layout_to_manifest(self, layout_data):
        """Minimal manifest from classify_elements() output (emergency fallback)."""
        elements = []
        for i, e in enumerate(layout_data.get('elements', [])):
            entry = dict(e)
            entry.setdefault('manifestId', f'e-{i:03d}')
            entry.setdefault('confidence', 0.85)
            entry.setdefault('flagReason', None)
            entry.setdefault('classificationSource', 'auto')
            entry.setdefault('skipRender', False)
            elements.append(entry)
        svgs = []
        for i, s in enumerate(layout_data.get('svgElements', [])):
            entry = dict(s)
            entry.setdefault('manifestId', f'svg-{i:03d}')
            entry['type'] = 'svg_image'
            entry.setdefault('skipRender', False)
            entry['confidence'] = 1.0
            entry['flagReason'] = None
            entry['classificationSource'] = 'auto'
            svgs.append(entry)
        return {
            'schemaVersion': 1,
            'slideIndex': 0,
            'sourceFile': '',
            'slideBg': layout_data.get('slideBg', '#0A0A0F'),
            'slideWidthPx': layout_data.get('slideWidth', 1280),
            'slideHeightPx': layout_data.get('slideHeight', 720),
            'complexity': 'low',
            'triageFlags': [],
            'elements': elements,
            'svgElements': svgs,
            'notes': '',
        }

    def _add_run(self, paragraph, text, run_data, max_emoji_pt=None):
        """Add a formatted run to a paragraph, splitting emoji from regular text (#34).

        Mixed runs like "⚠ Wrong naming" are split into separate PPTX runs:
        one with emoji font for the symbol, one with Calibri for the text.
        This prevents colored squares from font mismatch.

        Args:
            max_emoji_pt: If set, clamp emoji font size to this value (pts).
                Used for small icon elements where emoji renders larger than
                the bounding box (#38 systemic icon overlap).
        """
        tt = run_data.get('textTransform', '')
        if tt == 'uppercase':
            text = text.upper()
        decoded = decode_entities(text)

        # Split into emoji and non-emoji segments
        segments = _split_emoji(decoded)
        last_run = None
        for seg_text, is_emoji in segments:
            if not seg_text:
                continue
            run = paragraph.add_run()
            run.text = seg_text
            run.font.name = EMOJI_FONT if is_emoji else self.font
            fs = run_data.get('fontSize', 14)
            pt_size = round(max(6, min(52, fs * 0.75)), 1)
            # Clamp emoji to fit bounding box for small icon elements (#38)
            if is_emoji and max_emoji_pt and pt_size > max_emoji_pt:
                pt_size = max_emoji_pt
            run.font.size = Pt(pt_size)
            run.font.bold = is_bold(run_data.get('fontWeight'))
            run.font.italic = run_data.get('fontStyle') == 'italic'
            ls = run_data.get('letterSpacing')
            if ls and ls != 'normal':
                try:
                    px = float(ls) if isinstance(ls, (int, float)) else float(str(ls).replace('px', ''))
                    if abs(px) >= 0.1:
                        rPr = run._r.get_or_add_rPr()
                        rPr.set('spc', str(int(px * 75)))
                except (ValueError, TypeError):
                    pass
            color = parse_css_color(run_data.get('color'), (255, 255, 255))
            if color:
                run.font.color.rgb = color
            last_run = run
        return last_run

    def build_table(self, slide, table_node, slide_bg_rgb):
        """Build a native PPTX table from extracted table data."""
        rows = table_node.get('rows', [])
        if not rows:
            return None
        num_rows = len(rows)
        num_cols = max(len(row) for row in rows) if rows else 0
        if num_cols == 0:
            return None
        r = table_node['rect']
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))

        table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, w, h)
        table = table_shape.table

        # Disable default table styling that causes white banded rows (#16)
        tbl = table._tbl
        tblPr = tbl.find(qn('a:tblPr'))
        if tblPr is not None:
            tblPr.attrib.clear()
            tblPr.set('bandRow', '0')
            tblPr.set('bandCol', '0')
            tblPr.set('firstRow', '0')
            tblPr.set('lastRow', '0')
            tblPr.set('firstCol', '0')
            tblPr.set('lastCol', '0')
            for child in list(tblPr):
                if child.tag.endswith('}tblStyle'):
                    tblPr.remove(child)

        # Set column widths from first row cell widths
        if rows[0]:
            for ci, cell_data in enumerate(rows[0]):
                if ci < num_cols:
                    col_w = self.px2emu_x(cell_data['rect']['w'])
                    tbl.tblGrid[ci].set('w', str(col_w))

        fallback_bg = RGBColor(slide_bg_rgb[0], slide_bg_rgb[1], slide_bg_rgb[2])

        for ri, row in enumerate(rows):
            for ci, cell_data in enumerate(row):
                if ci >= num_cols or ri >= num_rows:
                    continue
                cell = table.cell(ri, ci)
                bg = parse_css_color(cell_data['styles'].get('backgroundColor'), slide_bg_rgb)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg or fallback_bg
                tf = cell.text_frame
                tf.word_wrap = True
                runs = cell_data.get('runs', [])
                first_run = True
                for run_data in runs:
                    if run_data.get('br'):
                        tf.add_paragraph()
                        continue
                    text = run_data.get('text', '').strip()
                    if not text:
                        continue
                    if first_run:
                        p = tf.paragraphs[0]
                        first_run = False
                    else:
                        p = tf.paragraphs[-1]
                    run = p.add_run()
                    run.text = text
                    run.font.name = self.font
                    run.font.size = Pt(round(max(6, min(52, run_data.get('fontSize', 12) * 0.75)), 1))
                    color = parse_css_color(run_data.get('color'), (255, 255, 255))
                    if color:
                        run.font.color.rgb = color
                    if is_bold(run_data.get('fontWeight')):
                        run.font.bold = True
        return table_shape

    def build_badge(self, slide, node, slide_bg_rgb):
        """Build a pill/circle badge shape."""
        r = node['rect']
        # Bounds check: skip badges outside the slide viewport (#17)
        if r['x'] > self.src_w or r['y'] > self.src_h:
            return None
        if r['x'] + r['w'] < 0 or r['y'] + r['h'] < 0:
            return None
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))
        st = node.get('styles', {})
        br = st.get('borderRadius', 0)
        if br >= min(r['w'], r['h']) / 2:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            try:
                shape.adjustments[0] = 0.5  # max rounding = pill
            except Exception:
                pass
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            try:
                shape.adjustments[0] = min(0.15, br / max(1, min(r['w'], r['h'])))
            except Exception:
                pass
        bg = parse_css_color(st.get('backgroundColor'), slide_bg_rgb)
        if bg:
            shape.fill.solid()
            shape.fill.fore_color.rgb = bg
        else:
            shape.fill.background()
        shape.line.fill.background()
        text = node.get('text', '').strip()
        if text:
            tf = shape.text_frame
            tf.word_wrap = False
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = text
            fs = st.get('fontSize', 12)
            pt_size = round(max(6, fs * 0.75), 1)
            # Clamp emoji in small badges to fit bounding box (#38)
            is_emoji_text = all(_is_emoji_char(ch) for ch in text if not ch.isspace())
            if is_emoji_text and r['h'] < 60:
                max_pt = max(6, r['h'] * 0.55)
                pt_size = min(pt_size, max_pt)
            run.font.name = EMOJI_FONT if is_emoji_text else self.font
            run.font.size = Pt(pt_size)
            color = parse_css_color(st.get('color'), (255, 255, 255))
            if color:
                run.font.color.rgb = color

    # -------------------------------------------------------------------------
    # Manifest-driven build path (model-driven triage pipeline)
    # -------------------------------------------------------------------------

    def build_slide_from_manifest(self, prs, manifest, screenshot_fn=None):
        """
        Build a PPTX slide from a fully-resolved build manifest.

        Args:
            prs: python-pptx Presentation object
            manifest: dict conforming to the build manifest schema (from prep)
            screenshot_fn: callable(svg_rect) -> png_path for SVG rendering

        Returns:
            The created slide object
        """
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # Background
        bg_color = parse_css_color(manifest.get('slideBg'))
        bg = slide.background
        bg.fill.solid()
        bg.fill.fore_color.rgb = bg_color or parse_css_color('#0A0A0F')
        slide_bg_rgb = (bg_color[0], bg_color[1], bg_color[2]) if bg_color else (10, 10, 15)

        elements = manifest.get('elements', [])
        svg_elements = manifest.get('svgElements', [])

        def _active(e, *types):
            return e.get('type') in types and not e.get('skipRender') and e.get('type') != 'skip'

        # Rendering order: shapes → accent_bars → svgs → tables → badges → richtexts
        shapes = sorted([e for e in elements if _active(e, 'shape')],
                        key=lambda e: e.get('depth', 0))
        accent_bars = [e for e in elements if _active(e, 'accent_bar')]
        tables = [e for e in elements if _active(e, 'table')]
        badges = [e for e in elements if _active(e, 'badge', 'circle')]
        richtexts = [e for e in elements if _active(e, 'richtext')]
        svg_items = [e for e in svg_elements if not e.get('skipRender')]

        for s in shapes:
            self._render_shape(slide, s, slide_bg_rgb)

        for a in accent_bars:
            self._render_accent_bar(slide, a, slide_bg_rgb)

        if screenshot_fn:
            for svg in svg_items:
                self._render_svg_image(slide, svg, screenshot_fn)

        for td in tables:
            try:
                self.build_table(slide, td, slide_bg_rgb)
            except Exception:
                pass

        for bd in badges:
            self._render_badge(slide, bd, slide_bg_rgb)

        for rt in richtexts:
            self._render_richtext(slide, rt)

        notes = manifest.get('notes')
        if notes:
            try:
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

        return slide

    def _eff_rect(self, element):
        """Return resolvedRect if present, else rect. resolvedRect wins always."""
        return element.get('resolvedRect') or element['rect']

    def _render_shape(self, slide, element, slide_bg_rgb):
        """Render a shape element from manifest onto slide."""
        r = self._eff_rect(element)
        st = element.get('styles', {})
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(r['h'], self.src_h - max(0, r['y'])))
        if w < 10000 and h < 8000:
            return  # allow thin decorative lines

        br = st.get('borderRadius', 0)
        if br > 5:
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            try:
                shape.adjustments[0] = min(0.12, br / min(r['w'], r['h']))
            except Exception:
                pass
        else:
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)

        fill_c = parse_css_color(st.get('backgroundColor'), slide_bg_rgb)
        if fill_c:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_c
        else:
            shape.fill.background()

        bg_raw = st.get('background', '')
        if bg_raw and 'linear-gradient' in str(bg_raw):
            angle, stops = _parse_css_gradient(bg_raw)
            if stops:
                _apply_gradient_fill(shape, angle, stops, slide_bg_rgb)

        border_c = parse_css_color(st.get('borderColor'), slide_bg_rgb)
        bw = st.get('borderWidth', 0)
        if border_c and bw > 0.3:
            shape.line.color.rgb = border_c
            shape.line.width = Pt(max(0.25, bw * 0.75))
        else:
            shape.line.fill.background()
        shape.text_frame.clear()

    def _render_accent_bar(self, slide, element, slide_bg_rgb):
        """Render a left-border accent bar (PATTERN-002)."""
        r = self._eff_rect(element)
        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(r['w'])
        h = self.px2emu_y(r['h'])
        fill_color_str = element.get('fill') or element.get('styles', {}).get('backgroundColor')
        fill_c = parse_css_color(fill_color_str, slide_bg_rgb)
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
        if fill_c:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_c
        shape.line.fill.background()
        shape.text_frame.clear()

    def _render_richtext(self, slide, element):
        """Render a richtext element from manifest onto slide."""
        r = self._eff_rect(element)
        runs = element.get('runs', [])
        if not runs:
            return
        full_text = ''.join(run.get('text', '') for run in runs).strip()
        if not full_text:
            return

        rotation = element.get('rotation', 0)
        align = element.get('styles', {}).get('textAlign', 'left')

        x = self.px2emu_x(max(0, r['x']))
        y = self.px2emu_y(max(0, r['y']))
        w = self.px2emu_x(min(r['w'], self.src_w - max(0, r['x'])))
        h = self.px2emu_y(min(max(r['h'], 14), self.src_h - max(0, r['y'])))

        if w < 5000 or h < 5000:
            return

        # For rotated text, swap dimensions and apply PPTX rotation (#27)
        if rotation and abs(rotation) >= 5:
            w, h = h, w

        tb = slide.shapes.add_textbox(x, y, w, h)
        if rotation and abs(rotation) >= 5:
            tb.rotation = rotation
        tb.fill.background()
        tb.line.fill.background()
        tf = tb.text_frame
        tf.word_wrap = True
        ws = element.get('styles', {}).get('whiteSpace', '')
        if ws == 'nowrap':
            tf.word_wrap = False
        self._apply_margins(tf, element.get('styles', {}))

        p = tf.paragraphs[0]
        p.alignment = pp_align(align)
        p.space_before = Pt(0)
        p.space_after = Pt(0)

        # Emoji font size cap for small icon elements (#38)
        max_emoji_pt = None
        if r['h'] < 60 and r['w'] < 80:
            all_text = ''.join(rd.get('text', '') for rd in runs).strip()
            if all_text and all(_is_emoji_char(ch) for ch in all_text if not ch.isspace()):
                max_emoji_pt = max(6, r['h'] * 0.55)

        for run_data in runs:
            text = run_data.get('text', '')
            if not text:
                continue
            if run_data.get('br'):
                p = tf.add_paragraph()
                p.alignment = pp_align(align)
                p.space_before = Pt(0)
                p.space_after = Pt(0)
                continue
            if '\n' in text:
                parts = text.split('\n')
                for pi, part in enumerate(parts):
                    if pi > 0:
                        p = tf.add_paragraph()
                        p.alignment = pp_align(align)
                        p.space_before = Pt(0)
                        p.space_after = Pt(0)
                    if part.strip():
                        self._add_run(p, part, run_data, max_emoji_pt=max_emoji_pt)
                continue
            self._add_run(p, text, run_data, max_emoji_pt=max_emoji_pt)

    def _render_badge(self, slide, element, slide_bg_rgb):
        """Render a badge/circle element from manifest (delegates to build_badge)."""
        # build_badge already handles bounds checks and shape creation
        self.build_badge(slide, element, slide_bg_rgb)

    def _render_svg_image(self, slide, element, screenshot_fn):
        """Render an SVG image element from manifest using screenshot callback."""
        r = self._eff_rect(element)
        if r['w'] < 30 or r['h'] < 30:
            return
        if element.get('lines', 0) < 3:
            return
        png_path = screenshot_fn(r)
        if png_path and os.path.exists(png_path):
            x = self.px2emu_x(max(0, r['x']))
            y = self.px2emu_y(max(0, r['y']))
            w = self.px2emu_x(r['w'])
            h = self.px2emu_y(r['h'])
            try:
                slide.shapes.add_picture(png_path, x, y, w, h)
            except Exception:
                pass
